# app.py

import streamlit as st
import streamlit.components.v1 as components
import pandas as pd
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from PIL import Image
import datetime
from uuid import uuid4
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

st.set_page_config(page_title="Location Uploader & PPTX Export", layout="wide")


# --- 追加：JSで画面幅をクエリにセット ---
components.html(
     """
    <script>
      const w = window.innerWidth;
      const url = new URL(window.location);
      // すでに同じ値ならリロードしない
      if (url.searchParams.get('screen_width') != w) {
        url.searchParams.set('screen_width', w);
        window.location.search = url.searchParams.toString();
      }
    </script>
    """,
    height=0,
)

# 画面幅に応じたプレビュー列数の設定（3段階レスポンシブ）
# 画面幅読み込み
screen_w = int(st.query_params.get("screen_width", ["0"])[0])

# screen_w が 0（まだ取得前）ならPC扱い
if screen_w == 0 or screen_w >= 1200:
    PREVIEW_COLS = 4
elif screen_w >= 768:
    PREVIEW_COLS = 4
else:
    PREVIEW_COLS = 4

PREVIEW_ROWS     = 2
PREVIEW_PER_PAGE = PREVIEW_COLS * PREVIEW_ROWS
PADDING          = 1

# 全体のCSS調整
st.markdown(
    f"""
    <style>
    img {{ width:100% !important; height:auto !important; }}
    .stCheckbox label {{ font-size: 0.9em; }}
    .stDownloadButton button {{ font-size: 0.9em; padding: 0.4em 0.8em; }}
    </style>
    """,
    unsafe_allow_html=True,
)

# --- 定義しておく ---
subcats = {
    "ハウススタジオ": ["和風","洋風","一軒家","マンション","アパート"],
    "オフィス":      ["執務室","会議室","ロビー"],
    "商業施設":      ["ショッピングモール","遊園地","温泉","水族館/動物園/植物園",
                     "博物館/美術館","映画館","ボーリング/ゲームセンター/ビリヤード場",
                     "商店街","ホテル"],
    "学校":          ["小学校","中学校","高校","大学/専門学校","幼稚園/保育園"],
    "病院":          ["受付","手術室"],
    "店舗":          ["コンビニ","ドラッグストア","スーパー","アパレル","ガソリンスタンド"],
    "飲食店":        ["中華料理屋","レストラン","カフェ","居酒屋","食堂","BAR"],
    "自然":          ["山","川","海","草原","森","湖/池","花畑","道"],
    "その他":        ["駐車場","屋上","神社仏閣","オープンスペース","夜景/イルミネーション",
                     "公民館","スポーツ施設"],
    "該当なし":      ["その他"]
}
detail_opts = [
    "ゼネ車使用可否","キッチン使用可否","同録の可否","養生の有無",
    "電源の有無","駐車場の有無","特機の使用可否","スモークの使用可否","火器の使用可否",
]


st.title("📍 ロケ地Info・資料出力アプリ")

# --- 1. 基本情報 ---
st.header("1. ロケ地基本情報")
loc_name       = st.text_input("ロケ地名")
address        = st.text_input("住所")
hp_link        = st.text_input("HPリンク（URL）")

# ロケ地種類：2カラム
c1, c2 = st.columns(2)
with c1:
    cat_main_val = st.selectbox("ロケ地種類（大分類）", list(subcats.keys()))
with c2:
    cat_sub_val  = st.selectbox("ロケ地種類（小分類）", subcats[cat_main_val])

transport_info = st.text_area("交通機関情報", height=80)

# 面積・天高：2カラム
a1, a2 = st.columns(2)
with a1:
    area_val   = st.number_input("面積 [m²]", min_value=0.0)
with a2:
    ceiling_val= st.number_input("天高 [cm]", min_value=0.0)

# 窓口連絡先：名前 + 電話3分割 + メール
st.subheader("窓口連絡先")
contact_person = st.text_input("窓口の担当者名")
ph1, ph2  = st.columns(2)
with ph1:
    phone1 = st.text_input("電話番号")
with ph2:
    contact_email = st.text_input("窓口のメールアドレス")

st.markdown("---")

# --- 2. 利用情報 --
st.header("2. 利用情報")
d1, d2 = st.columns(2)
with d1:
    price_day  = st.number_input("金額／day（¥）", min_value=0)
with d2:
    price_hour = st.number_input("金額／h（¥）",  min_value=0)

price_note = st.text_area("金額備考", height=80)

# 利用可能日時：24h + 開始・終了
st.subheader("利用可能日時")
t1, t2, t3 = st.columns([1,1,1])
with t1:
    open_24h   = st.checkbox("24時間利用可")
with t2:
    start_time = st.time_input("開始時間", value=datetime.time(9,0))
with t3:
    end_time   = st.time_input("終了時間", value=datetime.time(18,0))

st.markdown("---")

# --- 3. 詳細オプション ---
st.header("3. 詳細オプション")
detail_values = {}
for i in range(0, len(detail_opts), 3):
    cols = st.columns(3)
    for j, opt in enumerate(detail_opts[i:i+3]):
        detail_values[opt] = cols[j].selectbox(opt, ["あり","なし","不明"], index=2, key=opt)

# 支払い方法
payment    = st.selectbox("支払い方法", ["現金","カード","請求書","不明"], index=3)
pay_note   = st.text_input("支払い備考")
# 使用可能人数
st.subheader("使用可能人数")
u1, u2, u3, u4 = st.columns([1,1,1,1])
with u1:
    specify_num   = st.checkbox("人数指定")
with u2:
    max_number    = st.number_input("最大人数", min_value=0) if specify_num else None
with u3:
    unlimited     = st.checkbox("上限なし")
with u4:
    unknown_count = st.checkbox("不明")



st.markdown("---")

# --- 4. 対象作品 & 担当者 ---
st.header("4. 対象作品")
work_no     = st.text_input("作品番号")
pm, p, co   = st.columns(3)
with pm:
    pm_person   = st.text_input("担当者 PM")
with p:
    p_person    = st.text_input("担当者 P")
with co:
    coordinator = st.text_input("ロケコーディネーター")

st.markdown("---")

# --- 5. Excel 保存ボタン ---
st.header("5. 入力データをExcelで保存")
if st.button("💾 データをExcelでダウンロード"):
    # 辞書にまとめる
    data = {
        "ロケ地名":       loc_name,
        "住所":           address,
        "HPリンク":       hp_link,
        "大分類":         cat_main_val,
        "小分類":         cat_sub_val,
        "交通機関情報":   transport_info,
        "面積[m²]":       area_val,
        "天高[cm]":       ceiling_val,
        "担当者名":       contact_person,
        "電話番号":       phone1,
        "メールアドレス":  contact_email,
        "金額/day":       price_day,
        "金額/h":         price_hour,
        "金額備考":       price_note,
        "24時間可":       open_24h,
        "開始時間":       start_time.strftime("%H:%M"),
        "終了時間":       end_time.strftime("%H:%M"),
        **detail_values,
        "人数指定":       specify_num,
        "最大人数":       max_number,
        "上限なし":       unlimited,
        "不明(人数)":     unknown_count,
        "支払い方法":     payment,
        "支払い備考":     pay_note,
        "作品番号":       work_no,
        "担当者 PM":      pm_person,
        "担当者 P":       p_person,
        "コーディネーター": coordinator,
    }
    df = pd.DataFrame([data])
    buf = io.BytesIO()
    with pd.ExcelWriter(buf, engine="openpyxl") as writer:
        df.to_excel(writer, index=False, sheet_name="ロケ地情報")
    buf.seek(0)

    st.download_button(
        label="📥 Excelをダウンロード",
        data=buf,
        file_name="location_data.xlsx",
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    )

st.markdown("---")

# --- 6. 画像アップロード & PPTX 出力 ---
# --- Configuration ---

# Inputs for location metadata
location_name = loc_name
address = address

# --- ユーザー入力：1スライドあたりの画像枚数を選択 ---

total_per_slide = st.selectbox(
    "1スライドあたりの画像枚数",
    [6, 9],
    index=0
)
if total_per_slide == 6:
    PPT_COLS, PPT_ROWS = 3, 2
else:
    PPT_COLS, PPT_ROWS = 3, 3
PPT_PER_SLIDE = PPT_COLS * PPT_ROWS

# --- フォントサイズ定義 ---
HEADING_FONT = Pt(20)   # スライド上部の見出し
LOC_FONT     = Pt(20)   # ロケ地名
TABLE_FONT   = Pt(10)   # 表の文字

# --- プレビュー設定（変更なし）---



# --- カテゴリ定義 ---
# --- 6. 画像アップロード & Preview ---
categories = [
    ("サムネイル：1枚のみ", "thumbs", False),
    ("ロケ地写真", "photos", True),
    ("アングル写真", "angles", True),
    ("その他設備・搬入搬出経路", "others", True),
    ("平面図", "floor", True),
    ("ロケ地MAP", "map_img", True),
]

def display_image(img, **kwargs):
    """
    Streamlit のバージョン差異を吸収して画像表示。
    kwargs に use_container_width=True などを渡せば
    存在する方の引数で呼び出します。
    """
    try:
        st.image(img, **kwargs)
    except TypeError:
        # 新 API に use_container_width がない場合はこちら
        # kwargs の中に use_container_width があれば削除して再試行
        fallback = kwargs.copy()
        fallback.pop("use_container_width", None)
        st.image(img, **fallback)

# セッション初期化（省略）
for _, key, _ in categories:
    st.session_state.setdefault(f"{key}_data", {})
    st.session_state.setdefault(f"{key}_include", {})
    st.session_state.setdefault(f"{key}_page", 1)
    st.session_state.setdefault(f"{key}_ctr", 0)

for label, key, multi in categories:
    data = st.session_state[f"{key}_data"]
    ctr = st.session_state[f"{key}_ctr"]
    uploaded = st.file_uploader(label, type=["png","jpg","jpeg"],
                                accept_multiple_files=multi,
                                key=f"upl_{key}_{ctr}")
    files = uploaded if isinstance(uploaded, list) else ([uploaded] if uploaded else [])
    if files:
        for f in files:
            if f.name not in data:
                data[f.name] = Image.open(f)
                st.session_state[f"{key}_include"][f.name] = True
        st.session_state[f"{key}_ctr"] += 1
        st.rerun()

    items = list(data.items())
    if not items:
        continue

    # ページ計算
    total_pages = (len(items) + PREVIEW_PER_PAGE - 1) // PREVIEW_PER_PAGE
    page = st.session_state[f"{key}_page"]
    page = max(1, min(page, total_pages))
    st.session_state[f"{key}_page"] = page

    # プレビュー表示
    st.markdown(
        f"<div style='border:1px solid #ddd; padding:{PADDING}px; margin-bottom:{PADDING}px; max-height:400px; overflow-y:auto;'>",
        unsafe_allow_html=True,
    )

    # ── プレビュー画像のループ（削除をチェックボックスで） ──
    # ── プレビュー表示ループ（チェックボックスで「削除」と「資料出力」）──
    cols_ui = st.columns(PREVIEW_COLS)
    start = (page - 1) * PREVIEW_PER_PAGE

    for idx, (name, img) in enumerate(items[start:start + PREVIEW_PER_PAGE]):
        col = cols_ui[idx % PREVIEW_COLS]
        with col:
            display_image(img, use_container_width=True)

            # 「資料出力」のチェック
            inc = st.checkbox(
                "資料出力",
                key=f"inc_{key}_{name}",
                value=st.session_state[f"{key}_include"][name]
            )
            st.session_state[f"{key}_include"][name] = inc

            # 「削除」のチェックをオンにしたら即削除
            delete = st.checkbox(
                "削除",
                key=f"del_{key}_{name}"
            )
            if delete:
                data.pop(name)
                st.session_state[f"{key}_include"].pop(name, None)
                new_n     = len(data)
                new_total = max(1, (new_n + PREVIEW_PER_PAGE - 1) // PREVIEW_PER_PAGE)
                st.session_state[f"{key}_page"] = min(page, new_total)
                st.rerun()




    st.markdown("</div>", unsafe_allow_html=True)

    # ページナビ
    if total_pages > 1:
        sel = st.radio(
            f"ページ ({page}/{total_pages})",
            [str(i) for i in range(1, total_pages+1)],
            index=page-1,
            horizontal=True,
            key=f"nav_{key}"
        )
        new_page = int(sel)
        if new_page != page:
            st.session_state[f"{key}_page"] = new_page
            st.rerun()

# --- PPTX 生成＆ダウンロード ---
if st.button("💾 PPTX を生成"):
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 共通：Blank レイアウト取得
    try:
        blank_layout = prs.slide_layouts[6]
    except IndexError:
        blank_layout = next(
            (lay for lay in prs.slide_layouts if lay.name.lower()=="blank"),
            prs.slide_layouts[5]
        )

    # --- 1枚目：サムネイルスライド ---
    thumb_slide = prs.slides.add_slide(blank_layout)
    # 上部にロケ地名
    TITLE_W = Inches(10)
    left    = (prs.slide_width - TITLE_W) / 2
    tb = thumb_slide.shapes.add_textbox(left, Inches(0.2), TITLE_W, Inches(0.6))
    p  = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = location_name
    p.alignment  = PP_ALIGN.CENTER
    run.font.name = "YuGothic"
    run.font.size = Pt(24)

    # サムネイル画像を中央に縮小表示（幅を60%に）
    thumbs = list(st.session_state["thumbs_data"].values())
    if thumbs:
        img   = thumbs[0]
        pic_w = prs.slide_width * 0.6
        pic_h = pic_w * img.height / img.width
        left  = (prs.slide_width - pic_w) / 2
        top   = (prs.slide_height - pic_h) / 2 + Inches(0.2)
        buf = io.BytesIO()
        img.save(buf, format=img.format or "PNG")
        buf.seek(0)
        thumb_slide.shapes.add_picture(buf, left, top, width=pic_w, height=pic_h)

    # --- 2枚目：メタデータスライド ---
    meta_slide = prs.slides.add_slide(blank_layout)

    # (A) タイトル「ロケ地情報」
    META_TITLE_W = Inches(10)
    left_title   = (prs.slide_width - META_TITLE_W) / 2
    tb_meta_title = meta_slide.shapes.add_textbox(
        left_title, Inches(0.2), META_TITLE_W, Inches(0.6)
    )
    p_title = tb_meta_title.text_frame.paragraphs[0]
    run_title = p_title.add_run()
    run_title.text = "ロケ地情報"
    p_title.alignment = PP_ALIGN.CENTER
    run_title.font.name = "YuGothic"
    run_title.font.size = Pt(24)

    # すべてのフィールドをリストに
    fields = [
        ("ロケ地名",       location_name),
        ("住所",           address),
        ("HPリンク",       hp_link),
        ("大分類",         cat_main_val),
        ("小分類",         cat_sub_val),
        ("交通機関情報",   transport_info),
        ("面積[m²]",       area_val),
        ("天高[cm]",       ceiling_val),
        ("担当者名",       contact_person),
        ("電話番号",       phone1),
        ("メールアドレス",  contact_email),
        ("金額/day",       price_day),
        ("金額/h",         price_hour),
        ("金額備考",       price_note),
        ("24時間可",       open_24h),
        ("開始時間",       start_time.strftime("%H:%M")),
        ("終了時間",       end_time.strftime("%H:%M")),
        *[(k, detail_values[k]) for k in detail_values],
        ("人数指定",       specify_num),
        ("最大人数",       max_number),
        ("上限なし",       unlimited),
        ("不明(人数)",     unknown_count),
        ("支払い方法",     payment),
        ("支払い備考",     pay_note),
        ("作品番号",       work_no),
        ("担当者 PM",      pm_person),
        ("担当者 P",       p_person),
        ("コーディネーター", coordinator),
    ]

    # 2分割
    mid = len(fields) // 2
    left_fields  = fields[:mid]
    right_fields = fields[mid:]

    # 各コラムの基本設定
    #margin_x  = Inches(0.7)
    #half_w    = (prs.slide_width - margin_x*2) / 2
    #label_w   = Inches(2.5)   # 左コラムはメールアドレスに合わせて広め
    #value_w   = half_w - label_w - Inches(0.2)
    #row_h     = Inches(0.25)  # 行間
    #start_y   = Inches(0.5)

    # ↓↓ ここから調整可能 ↓↓
    # 項目のY開始位置を0.2→1.0インチに下げてタイトルと被らなくする
    start_y   = Inches(1.0)

    # 左右マージン
    margin_x  = Inches(0.7)
    # コラム幅（左右で均等）
    half_w    = (prs.slide_width - margin_x*2) / 2
    # 左コラム：ラベル幅を2.5→2.2インチに少し狭く
    label_w   = Inches(2.0)
    # 値幅は自動計算
    value_w   = half_w - label_w - Inches(0.1)
    # 行の高さ（行間）は0.25→0.3インチに調整
    row_h     = Inches(0.3)
    # ↑↑ ここまで調整可能 ↑↑


    # 左コラム
    for i, (label, val) in enumerate(left_fields):
        y = start_y + row_h * i
        # ラベル
        tb_lab = meta_slide.shapes.add_textbox(margin_x, y, label_w, row_h)
        p_lab  = tb_lab.text_frame.paragraphs[0]
        p_lab.text = label
        p_lab.font.name = "YuGothic"
        p_lab.font.size = Pt(10)
        # 値
        tb_val = meta_slide.shapes.add_textbox(
            margin_x + label_w + Inches(0.1), y, value_w, row_h
        )
        p_val  = tb_val.text_frame.paragraphs[0]
        p_val.text = str(val)
        p_val.font.name = "YuGothic"
        p_val.font.size = Pt(10)

    # 右コラム（Coordinator に合わせてラベル幅設定）
    label_w_r = Inches(1.8)
    value_w_r = half_w - label_w_r - Inches(0.2)
    for i, (label, val) in enumerate(right_fields):
        y = start_y + row_h * i
        x0 = margin_x + half_w
        tb_lab = meta_slide.shapes.add_textbox(x0, y, label_w_r, row_h)
        p_lab  = tb_lab.text_frame.paragraphs[0]
        p_lab.text = f"{label}"
        p_lab.font.name = "YuGothic"
        p_lab.font.size = Pt(10)
        tb_val = meta_slide.shapes.add_textbox(x0 + label_w_r + Inches(0.1), y, value_w_r, row_h)
        p_val  = tb_val.text_frame.paragraphs[0]
        p_val.text = f"{val}"
        p_val.font.name = "YuGothic"
        p_val.font.size = Pt(10)

    # --- 3枚目以降：その他カテゴリの画像スライド（省略せず従来どおり） ---
    for label, key, _ in categories:
        if key == "thumbs":
            continue
        imgs = [
            img for name,img in st.session_state[f"{key}_data"].items()
            if st.session_state[f"{key}_include"][name]
        ]
        if not imgs:
            continue

        for i in range(0, len(imgs), PPT_PER_SLIDE):
            chunk = imgs[i:i + PPT_PER_SLIDE]
            slide = prs.slides.add_slide(blank_layout)

            # カテゴリ見出し
            tb_cat = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.3),
                Inches(3), Inches(0.5)
            )
            p_cat    = tb_cat.text_frame.paragraphs[0]
            run_cat  = p_cat.add_run()
            run_cat.text = label
            p_cat.alignment = PP_ALIGN.LEFT
            run_cat.font.name = "YuGothic"
            run_cat.font.size = HEADING_FONT

            # ロケ地名
            TEXT_W = Inches(10)
            LEFT   = (prs.slide_width - TEXT_W) / 2
            tb_loc2 = slide.shapes.add_textbox(LEFT, Inches(0.3), TEXT_W, Inches(0.5))
            p_loc2  = tb_loc2.text_frame.paragraphs[0]
            run_loc2= p_loc2.add_run()
            run_loc2.text = location_name
            p_loc2.alignment = PP_ALIGN.CENTER
            run_loc2.font.name = "YuGothic"
            run_loc2.font.size = HEADING_FONT

            # 画像グリッド…
            usable_w = prs.slide_width - Inches(1)
            usable_h = prs.slide_height - Inches(1.5)
            gap_w, gap_h = Inches(0.2), Inches(0.2)
            cell_w = (usable_w - gap_w*(PPT_COLS-1)) / PPT_COLS
            cell_h = (usable_h - gap_h*(PPT_ROWS-1)) / PPT_ROWS
            left_m, top_m = Inches(0.5), Inches(1.5)

            for idx, img in enumerate(chunk):
                r, c = divmod(idx, PPT_COLS)
                x = left_m + c*(cell_w+gap_w)
                y = top_m + r*(cell_h+gap_h)
                ow, oh = img.size
                ratio, cell_ratio = ow/oh, cell_w/cell_h
                if ratio > cell_ratio:
                    pw, ph = cell_w, cell_w/ratio
                else:
                    ph, pw = cell_h, cell_h*ratio
                px = x + (cell_w-pw)/2
                py = y + (cell_h-ph)/2
                #if pw<cell_w or ph<cell_h:
                #    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cell_w, cell_h)
                #    bg.fill.solid()
                #    bg.fill.fore_color.rgb = RGBColor(255,255,255)
                #    bg.line.fill.background()
                buf = io.BytesIO()
                img.save(buf, format=img.format or "PNG")
                buf.seek(0)
                slide.shapes.add_picture(buf, px, py, width=pw, height=ph)

    # 保存＆ストア
    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    st.session_state['pptx_bytes'] = out.getvalue()

# --- ダウンロードボタン（ifブロック外） ---
if 'pptx_bytes' in st.session_state:
    st.download_button(
        "📥 PPTXをダウンロード",
        data=st.session_state['pptx_bytes'],
        file_name="location_pictures.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
